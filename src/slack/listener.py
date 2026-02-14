"""
NEXUS Slack Listener (v1)

Listens for messages in #garrett-nexus and routes them directly to the
reasoning engine. No more HTTP round-trip to the server.

Keeps: file download/parsing, markdown-to-slack conversion.
Removes: old server routing, category formatting, in-memory history.
"""

import asyncio
import logging
import os
import re
import tempfile

import aiohttp

logger = logging.getLogger("nexus.slack.listener")
from slack_sdk.socket_mode.aiohttp import SocketModeClient
from slack_sdk.socket_mode.request import SocketModeRequest
from slack_sdk.socket_mode.response import SocketModeResponse
from slack_sdk.web.async_client import AsyncWebClient

from src.agents.haiku_intake import format_with_tool_result, run_haiku_intake
from src.agents.intake_dispatcher import dispatch
from src.agents.planner import create_plan
from src.agents.registry import registry
from src.config import get_key
from src.documents.generator import generate_document
from src.memory.store import memory
from src.ml.rag import build_rag_context, ingest_conversation
from src.ml.similarity import analyze_new_directive
from src.observability.logging import thread_ts_var
from src.sessions.cli_pool import cli_pool


def md_to_slack(text: str) -> str:
    """Convert markdown/HTML to Slack mrkdwn format.

    Handles code blocks, inline code, bold, links, lists, and HTML entities.
    Preserves code blocks verbatim (Slack renders ``` natively).
    """
    # Step 1: Extract code blocks to protect them from other transformations
    code_blocks = []
    def _save_code_block(match):
        code_blocks.append(match.group(0))
        return f"\x00CODE_BLOCK_{len(code_blocks) - 1}\x00"

    # Match fenced code blocks (```lang\n...\n```) — strip language hint
    text = re.sub(r'```\w*\n(.*?)```', lambda m: f"```\n{m.group(1)}```", text, flags=re.DOTALL)
    text = re.sub(r'```.*?```', _save_code_block, text, flags=re.DOTALL)

    # Step 2: Strip HTML tags (Claude Code CLI can output HTML)
    text = re.sub(r'<(?!http|mailto|#|@|!)/?[a-zA-Z][^>]*>', '', text)

    # Step 3: HTML entities
    text = text.replace('&amp;', '&')
    text = text.replace('&lt;', '<')
    text = text.replace('&gt;', '>')
    text = text.replace('&quot;', '"')
    text = text.replace('&#39;', "'")
    text = text.replace('&nbsp;', ' ')

    # Step 4: Markdown → Slack mrkdwn
    text = re.sub(r'^#{1,6}\s+(.+)$', r'*\1*', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.+?)\*\*', r'*\1*', text)
    text = re.sub(r'__(.+?)__', r'*\1*', text)
    text = re.sub(r'(?<![`*])`([^`\n]+)`(?![`*])', r'`\1`', text)  # inline code
    text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<\2|\1>', text)
    text = re.sub(r'^[\-\*]\s+', '•  ', text, flags=re.MULTILINE)
    text = re.sub(r'^\d+\.\s+', '•  ', text, flags=re.MULTILINE)  # numbered lists
    text = re.sub(r'^[\-\*]{3,}$', '———', text, flags=re.MULTILINE)
    text = text.replace('***', '*')

    # Step 5: Clean up excessive blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Step 6: Restore code blocks
    for i, block in enumerate(code_blocks):
        text = text.replace(f"\x00CODE_BLOCK_{i}\x00", block)

    return text.strip()


def format_code_output(output: str, agent_name: str = "") -> list[dict]:
    """Format agent output with code into Slack Block Kit blocks.

    Returns a list of Slack blocks. Code sections get their own code blocks.
    Non-code text gets mrkdwn formatting.
    """
    blocks = []

    if agent_name:
        blocks.append({
            "type": "context",
            "elements": [{"type": "mrkdwn", "text": f"*{agent_name}* completed:"}],
        })

    # Split on code fences
    parts = re.split(r'(```\w*\n.*?```)', output, flags=re.DOTALL)

    for part in parts:
        part = part.strip()
        if not part:
            continue

        if part.startswith('```'):
            # Extract code content (strip fences and language hint)
            code = re.sub(r'^```\w*\n?', '', part)
            code = re.sub(r'\n?```$', '', code)
            if len(code) > 2900:
                code = code[:2900] + "\n... (truncated)"
            blocks.append({
                "type": "section",
                "text": {"type": "mrkdwn", "text": f"```{code}```"},  # type: ignore[dict-item]
            })
        else:
            text = md_to_slack(part)
            if len(text) > 2900:
                text = text[:2900] + "\n..."
            if text:
                blocks.append({
                    "type": "section",
                    "text": {"type": "mrkdwn", "text": text},  # type: ignore[dict-item]
                })

    return blocks


async def download_and_parse_file(url: str, filename: str, bot_token: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, headers={"Authorization": f"Bearer {bot_token}"}) as resp:
                if resp.status != 200:
                    return f"(Failed to download: HTTP {resp.status})"

                text_exts = {"py", "js", "ts", "jsx", "tsx", "json", "yaml", "yml",
                             "html", "css", "md", "csv", "txt", "xml", "sh", "bash",
                             "sql", "toml", "ini", "cfg", "env", "log"}
                if ext in text_exts:
                    content = await resp.text()
                    if len(content) > 5000:
                        content = content[:5000] + "\n... (truncated)"
                    return f"```\n{content}\n```"

                file_bytes = await resp.read()
                with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name

                try:
                    content = ""
                    if ext in ("docx", "doc"):
                        from docx import Document
                        doc = Document(tmp_path)
                        content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                    elif ext in ("xlsx", "xls"):
                        import openpyxl
                        wb = openpyxl.load_workbook(tmp_path, read_only=True)
                        rows = []
                        for sheet in wb.sheetnames[:3]:
                            ws = wb[sheet]
                            rows.append(f"--- Sheet: {sheet} ---")
                            for row in ws.iter_rows(max_row=50, values_only=True):
                                rows.append(" | ".join(str(c) if c else "" for c in row))
                        content = "\n".join(rows)
                    elif ext in ("pptx", "ppt"):
                        from pptx import Presentation
                        prs = Presentation(tmp_path)
                        slides = []
                        for i, slide in enumerate(prs.slides, 1):
                            texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                            if texts:
                                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
                        content = "\n\n".join(slides)
                    elif ext == "pdf":
                        try:
                            import fitz
                            pdf_doc = fitz.open(tmp_path)
                            content = "\n".join(page.get_text() for page in pdf_doc[:10])
                        except ImportError:
                            content = "(Install PyMuPDF for PDF reading)"
                    elif ext in ("png", "jpg", "jpeg", "gif", "webp", "svg", "bmp", "ico"):
                        # Save image and return path for CLI to read with Read tool
                        import shutil
                        img_dir = os.path.join(tempfile.gettempdir(), "nexus_images")
                        os.makedirs(img_dir, exist_ok=True)
                        # Sanitize filename to prevent path traversal attacks
                        safe_filename = os.path.basename(filename)
                        if not safe_filename:
                            safe_filename = "image"
                        img_path = os.path.join(img_dir, safe_filename)
                        shutil.copy2(tmp_path, img_path)
                        content = f"\n[IMAGE ATTACHED: Use your Read tool to view the image at {img_path}]\n"
                    else:
                        content = f"(Unsupported file type: .{ext})"

                    if len(content) > 8000:
                        content = content[:8000] + "\n... (truncated)"
                    return f"```\n{content}\n```" if content else "(Empty or unreadable)"
                finally:
                    os.unlink(tmp_path)

    except Exception as e:
        return f"(Failed to read {filename}: {e})"


async def upload_file_to_slack(web_client: AsyncWebClient, channel: str,
                                filepath: str, title: str = "", comment: str = "",
                                thread_ts: str | None = None):
    try:
        kwargs = dict(
            channel=channel, file=filepath,
            title=title or os.path.basename(filepath), initial_comment=comment,
        )
        if thread_ts:
            kwargs["thread_ts"] = thread_ts
        await web_client.files_upload_v2(**kwargs)  # type: ignore[arg-type]
    except Exception as e:
        logger.error("File upload failed: %s", e)


_web_client: AsyncWebClient | None = None
_channel_id: str | None = None


def get_slack_client() -> AsyncWebClient | None:
    return _web_client


def get_channel_id() -> str | None:
    return _channel_id


def _build_execution_report(cli_result, directive_text: str) -> str:
    """Build a concise post-execution summary report for Slack."""
    from src.agents.task_result import TaskResult

    if not isinstance(cli_result, TaskResult):
        return ""

    parts: list[str] = []
    parts.append(":clipboard: *Execution Report*")

    # Status
    status_icon = ":white_check_mark:" if cli_result.succeeded else ":x:"
    parts.append(f"{status_icon} *Status:* {cli_result.status}")

    # Duration
    if cli_result.elapsed_seconds > 0:
        mins, secs = divmod(int(cli_result.elapsed_seconds), 60)
        time_str = f"{mins}m {secs}s" if mins else f"{secs}s"
        parts.append(f":stopwatch: *Duration:* {time_str}")

    # Cost
    if cli_result.cost_usd > 0:
        parts.append(f":moneybag: *CLI Cost:* ${cli_result.cost_usd:.4f}")

    # Tools used
    meta = cli_result.metadata or {}
    tools_used = meta.get("tools_used", 0)
    if tools_used:
        # Summarize tool usage by type
        tools_log: list[str] = meta.get("tools_log", [])
        if tools_log:
            from collections import Counter
            counts = Counter(tools_log)
            tool_summary = ", ".join(
                f"{name} x{count}" for name, count in counts.most_common(8)
            )
            parts.append(f":hammer_and_wrench: *Tools:* {tools_used} total ({tool_summary})")
        else:
            parts.append(f":hammer_and_wrench: *Tools:* {tools_used} total")

    # Files touched
    files_touched: list[str] = meta.get("files_touched", [])
    if files_touched:
        file_names = [f.split("/")[-1] for f in files_touched[:10]]
        parts.append(f":page_facing_up: *Files:* {', '.join(file_names)}")
        if len(files_touched) > 10:
            parts.append(f"  _...and {len(files_touched) - 10} more_")

    # Errors
    if not cli_result.succeeded:
        if cli_result.error_type:
            parts.append(f":warning: *Error type:* {cli_result.error_type}")
        if cli_result.error_detail:
            parts.append(f":memo: *Detail:* {cli_result.error_detail[:200]}")

    return "\n".join(parts)


def _is_garbage_response(response: str) -> bool:
    """Detect responses that are leaked instructions or placeholder templates.

    When CLI exits prematurely, it can emit raw system prompt fragments or
    template placeholders instead of real answers.
    """
    r = response.strip().lower()

    # Short responses are more likely to be garbage — real answers are usually >50 chars
    is_short = len(r) < 200

    # Exact substring patterns that always indicate garbage
    garbage_patterns = [
        "e.g.,",
        "e.g.)",
        "(e.g.",
        "e.g.'",
        "completed x%",
        "currently on section y",
        "not started yet'",
        "provide current status",
        "replace this with",
        "[insert",
        "[your ",
        "[fill in",
        "[placeholder",
        "fill in the blank",
        "your response here",
        "write your answer",
        "template instructions",
        "placeholder text",
        "example placeholders",
    ]
    for pattern in garbage_patterns:
        if pattern in r:
            return True

    # Detect bare instruction verbs followed by template-style content
    # Covers both double and single quotes, parens with examples
    if re.match(
        r'^(provide|list|describe|explain|summarize|generate|create|write|include|add|state|give|offer)\s'
        r".{0,60}(e\.g\.|for example|such as|\(.*?\)|\".*?\"|'.*?')",
        r,
    ):
        return True

    # Detect responses that are entirely a single instruction sentence with
    # parenthetical or quoted examples — the hallmark of a leaked system prompt
    if re.match(
        r'^(provide|list|describe|explain|summarize|state|give|offer|respond|answer)\s.{10,150}(\(.*?\)|\'.*?\'|".*?")\s*$',
        r,
    ):
        return True

    # Short responses starting with an imperative verb are suspicious
    if is_short and re.match(
        r'^(provide|list|describe|explain|summarize|state|give|offer|respond with|answer with)\s',
        r,
    ):
        return True

    # Detect responses that quote placeholder values like 'X%', 'section Y'
    return bool(
        re.search(r"['\"]\w?%['\"]\s*[,.]", r)
        or re.search(r"['\"]section \w['\"]", r)
    )


async def _fetch_thread_history(
    web_client: AsyncWebClient, channel: str, thread_ts: str, bot_user_id: str,
    current_ts: str, limit: int = 10,
) -> list[dict[str, str]]:
    """Fetch recent thread messages to build conversational context for CLI.

    Returns a list of {"role": "user"|"assistant", "text": "..."} dicts,
    excluding the current message (identified by current_ts).
    """
    try:
        result = await web_client.conversations_replies(
            channel=channel, ts=thread_ts, limit=limit + 5,
        )
        messages: list[dict] = result.get("messages", [])
        history: list[dict[str, str]] = []
        for msg in messages:
            msg_ts = msg.get("ts", "")
            if msg_ts == current_ts:
                continue
            # Skip the thread parent if it's a file upload with no text
            msg_text = msg.get("text", "").strip()
            if not msg_text:
                continue
            # Skip processing indicators
            if msg_text.startswith(":hourglass"):
                continue
            role = "assistant" if msg.get("user") == bot_user_id or msg.get("bot_id") else "user"
            history.append({"role": role, "text": msg_text[:2000]})
        # Return only the most recent `limit` messages
        return history[-limit:]
    except Exception as e:
        logger.warning("Thread history fetch failed (non-fatal): %s", e)
        return []


_SYSTEM_PREAMBLE = (
    "You are NEXUS, a 26-agent autonomous software engineering organization "
    "built by Garrett Eaglin. You are responding in the #garrett-nexus Slack channel.\n\n"
    "YOU ARE RUNNING IN THE NEXUS PROJECT DIRECTORY. You have full access to the codebase. "
    "Key files you can read:\n"
    "- docs/ARCHITECTURE.md — full architecture spec (v1.3, 29 sections)\n"
    "- README.md — project overview, ML system, usage examples\n"
    "- ORG_CHART.md — organizational structure\n"
    "- config/agents.yaml — all agent definitions\n"
    "- src/ — full Python source code\n"
    "- src/ml/ — ML self-learning system (router, predictor, embeddings, feedback)\n"
    "- src/orchestrator/ — LangGraph engine, task runner\n"
    "- src/slack/ — Slack listener and notifier\n"
    "- src/agents/ — agent base, SDK bridge, CEO interpreter\n\n"
    "RULES:\n"
    "1. When asked to create documents, READ the actual codebase first and include REAL data.\n"
    "2. When asked to enrich or fetch from GitHub, use your tools to actually read the repo files.\n"
    "3. Never output placeholder text, templates, or plans to do something — actually DO it.\n"
    "4. Never output instruction-style text like 'Provide X' or 'List Y' — give the actual answer.\n"
    "5. If asked about architecture, read docs/ARCHITECTURE.md and src/ to give real details.\n"
    "6. If you don't know something, say so directly.\n"
    "7. Always give complete, substantive answers grounded in the actual codebase.\n"
    "8. NEVER give the user step-by-step instructions to follow. YOU execute the steps yourself.\n"
    "   Do NOT say 'Step 1: Create a file...', 'Run npm install...', or 'Add this to your config'.\n"
    "   Instead, create the files, run the commands, and install the dependencies yourself.\n"
    "   Garrett is the CEO — he tells you WHAT to build, you figure out HOW and DO it.\n\n"
    "PROJECT SETUP — when creating ANY new project:\n"
    "1. Always run `git init` in the new project directory.\n"
    "2. Make an initial commit with all scaffolding files.\n"
    "3. Deploy to GitHub: run `gh repo create Garrett-s-Apps/<project-name> --public --source . --push`\n"
    "4. Install ALL dependencies your code imports before committing (npm install, pip install, etc.).\n"
    "5. Verify the project builds/runs before reporting completion.\n"
    "Never skip git initialization or GitHub deployment — Garrett expects every project on GitHub."
)

_RETRY_REINFORCEMENT = (
    "\n\nCRITICAL: Your previous attempt produced a garbage or placeholder response. "
    "You MUST respond with a real, substantive answer this time. Do NOT output "
    "template text, instructions, or example placeholders. Do NOT start with "
    "imperative verbs like 'Provide' or 'Describe' — those are instructions, "
    "not answers. Answer the user's actual question directly."
)


def _build_ml_briefing(message: str) -> str:
    """Build an ML intelligence briefing for the CLI prompt.

    Combines three ML signals: similar past directives, cost/quality predictions,
    and escalation risk. Returns empty string if no ML data is available yet.
    """
    try:
        analysis = analyze_new_directive(message)

        parts = []
        # Similar past work
        if analysis.get("has_precedent"):
            similar = analysis["similar_directives"][:3]
            items = []
            for s in similar:
                pct = int(s["similarity"] * 100)
                cost = f"${s['total_cost']:.2f}" if s["total_cost"] > 0 else "n/a"
                items.append(f"  - [{pct}% match] {s['directive_text'][:80]} (cost: {cost})")
            parts.append("Similar past work:\n" + "\n".join(items))

        # Cost estimate
        cost_est = analysis.get("cost_estimate", {})
        if cost_est.get("predicted"):
            parts.append(
                f"Cost estimate: ${cost_est['predicted']:.2f} "
                f"(range: ${cost_est['confidence_low']:.2f}–${cost_est['confidence_high']:.2f})"
            )

        # Historical averages
        hist = analysis.get("historical_average", {})
        if hist.get("avg_tasks", 0) > 0:
            parts.append(
                f"Historical average: {hist['avg_tasks']} tasks, "
                f"${hist['avg_cost']:.2f}, {hist['avg_duration_sec']/60:.0f}min"
            )

        # Risk assessment
        if analysis.get("risk_factors"):
            parts.append(f"Risk: {analysis['risk'].upper()} — " + "; ".join(analysis["risk_factors"]))

        # Agent recommendations
        recs = analysis.get("agent_recommendations", {})
        if recs:
            top = list(recs.items())[:3]
            parts.append("Top agents: " + ", ".join(f"{a} ({r:.0%})" for a, r in top))

        if not parts:
            return ""

        return (
            "\n\n[ML INTELLIGENCE BRIEFING — internal signals from NEXUS learning system]\n"
            + "\n".join(parts)
            + "\nUse this data to calibrate effort, cost awareness, and agent selection. "
            "Do not surface these metrics to the user unless they specifically ask about costs or history."
        )
    except Exception:
        return ""


def _build_threaded_prompt(
    history: list[dict[str, str]], current_message: str, *, is_retry: bool = False,
    ml_briefing: str = "", rag_context: str = "", execution_plan: str = "",
) -> str:
    """Build a prompt with thread history so CLI has conversational context."""
    parts = [_SYSTEM_PREAMBLE, ""]

    if ml_briefing:
        parts.append(ml_briefing)
        parts.append("")

    # Execution plan from Sonnet planner — gives Opus structured direction
    if execution_plan:
        parts.append(
            "[EXECUTION PLAN — created by Sonnet planner. Follow this plan unless "
            "you discover a better approach during implementation.]\n"
        )
        parts.append(execution_plan)
        parts.append("")

    # Thread history goes before RAG so the live conversation is closer to
    # the user message (LLMs attend more strongly to nearby context).
    if history:
        parts.append("You are continuing a Slack thread conversation. Here is the prior context:")
        parts.append("")
        for msg in history:
            prefix = "User" if msg["role"] == "user" else "You (NEXUS)"
            parts.append(f"{prefix}: {msg['text']}")
        parts.append("")

    # RAG context is supplementary reference — placed after thread history
    if rag_context:
        parts.append(rag_context)
        parts.append("")

    if history:
        parts.append(f"User's new message: {current_message}")
    else:
        parts.append(f"User's message: {current_message}")

    parts.append("")
    parts.append(
        "Respond to the user's message using the thread context above if present. "
        "Give a complete, substantive answer — do not just acknowledge the request."
    )

    if is_retry:
        parts.append(_RETRY_REINFORCEMENT)

    return "\n".join(parts)


async def start_slack_listener():
    global _web_client, _channel_id

    bot_token = get_key("SLACK_BOT_TOKEN")
    app_token = get_key("SLACK_APP_TOKEN")

    if not bot_token or not app_token:
        logger.warning("Missing tokens — Slack disabled")
        return

    web_client = AsyncWebClient(token=bot_token)
    socket_client = SocketModeClient(app_token=app_token, web_client=web_client)
    _web_client = web_client

    try:
        result = await web_client.conversations_list(types="public_channel,private_channel", limit=200)
        for ch in result["channels"]:
            if ch["name"] == "garrett-nexus":
                _channel_id = ch["id"]
                break
        if not _channel_id:
            logger.warning("Channel #garrett-nexus not found")
    except Exception as e:
        logger.error("Channel lookup failed: %s", e)

    auth = await web_client.auth_test()
    bot_user_id = auth["user_id"]

    async def handle_event(client: SocketModeClient, req: SocketModeRequest):
        await client.send_socket_mode_response(SocketModeResponse(envelope_id=req.envelope_id))

        if req.type != "events_api":
            return
        event = req.payload.get("event", {})
        subtype = event.get("subtype", "")
        # Allow file_share (images/attachments) through; block bot messages, edits, etc.
        ignored_subtypes = {"bot_message", "message_changed", "message_deleted", "channel_join", "channel_leave"}
        if event.get("type") != "message" or subtype in ignored_subtypes or event.get("user") == bot_user_id:
            return

        channel_id = event.get("channel", "")
        text = event.get("text", "").strip()

        files = event.get("files", [])
        file_context = ""
        if files:
            for f in files:
                fname = f.get("name", "unknown")
                url = f.get("url_private", "")
                if url:
                    file_context += f"\n[File: {fname}]\n"
                    file_context += await download_and_parse_file(url, fname, bot_token)
            text = f"{text}\n{file_context}" if text else f"[Sent files]{file_context}"

        if not text:
            return

        # Idempotency: skip duplicate messages (Slack Socket Mode at-least-once delivery)
        slack_ts = event.get("ts", "")
        client_msg_id = event.get("client_msg_id", "")
        dedup_key = f"{channel_id}:{slack_ts}:{client_msg_id or ''}"
        if memory.is_message_processed(dedup_key):
            logger.debug("Duplicate message detected, skipping: %s", dedup_key)
            return

        logger.info("Received: %s", text[:100])

        thinking_msg = None
        thread_ts = event.get("thread_ts") or event.get("ts")
        thread_ts_var.set(thread_ts or "")
        memory.emit_event("slack", "message_received", {
            "text": text[:200], "thread_ts": thread_ts,
            "user": event.get("user", ""), "has_files": bool(files),
        })

        try:
            async def safe_react(name):
                try:
                    await web_client.reactions_add(channel=channel_id, name=name, timestamp=event.get("ts"))
                except Exception:
                    pass  # reactions:write scope may not be granted

            # Send a "working on it" indicator so the user knows we received it
            try:
                result = await web_client.chat_postMessage(
                    channel=channel_id,
                    text=":hourglass_flowing_sand: Processing...",
                    thread_ts=thread_ts,
                )
                thinking_msg = result.get("ts")
            except Exception:
                pass

            # Get org context for Haiku intake
            org_summary = registry.get_org_summary()
            active = memory.get_active_directive()
            status_brief = f"Active directive: {active['text'][:100]}" if active else "Idle — no active directives"

            # Fetch thread history BEFORE intake so Haiku has conversation context
            thread_history = await _fetch_thread_history(
                web_client, channel_id, thread_ts, bot_user_id,
                current_ts=event.get("ts", ""),
            )

            # Detect follow-up: if this thread has an active CLI session,
            # tell Haiku this is a follow-up so it routes to start_directive
            thread_context = ""
            has_active_session = cli_pool.has_busy_sessions(thread_ts)
            if has_active_session:
                thread_context = (
                    "This thread has an ACTIVE CLI session — the user is following up on "
                    "engineering work already in progress. Route follow-up requests to "
                    "start_directive so the existing session can handle them. Only use "
                    "query tools if the user is asking a pure data question."
                )
            elif thread_history and len(thread_history) > 1:
                thread_context = (
                    "This is a continuing thread conversation. If the user's message "
                    "relates to prior engineering work discussed in the thread, use "
                    "start_directive to continue that work."
                )

            # Run Haiku intake to classify intent and route appropriately
            intake_result = await run_haiku_intake(
                message=text,
                thread_history=thread_history,
                org_summary=org_summary,
                system_status_brief=status_brief,
                thread_context=thread_context,
            )

            # Emit Haiku classification event for dashboard live ticker
            memory.emit_event("haiku", "classification", {
                "tool_called": intake_result.tool_called or "conversation",
                "tokens_in": intake_result.tokens_in,
                "tokens_out": intake_result.tokens_out,
                "message_preview": text[:100],
                "thread_ts": thread_ts,
            })

            response = None
            project_path = os.environ.get("NEXUS_PROJECT_PATH", os.path.expanduser("~/Projects/nexus"))

            # Handle based on tool called
            if intake_result.tool_called is None:
                # Pure conversation — Haiku already generated the response
                response = intake_result.response_text

            elif intake_result.tool_called == "start_directive":
                # Engineering work — use the existing CLI flow
                directive_text = intake_result.tool_input.get("directive_text", text) if intake_result.tool_input else text

                # Detect explicit cancel intent
                cancel_keywords = {"stop", "cancel", "nevermind", "never mind", "abort", "kill it"}
                is_cancel = any(kw in text.lower() for kw in cancel_keywords)

                if is_cancel and cli_pool.has_busy_sessions(thread_ts):
                    killed = await cli_pool.cancel_all(thread_ts)
                    memory.emit_event("cli", "cancelled_by_user", {
                        "thread_ts": thread_ts,
                        "sessions_killed": killed,
                        "cancel_message": text[:200],
                    })
                    response = f":octagonal_sign: *Cancelled {killed} running agent(s).* Say what you'd like to do next."
                    if thinking_msg:
                        try:
                            await web_client.chat_delete(channel=channel_id, ts=thinking_msg)
                        except Exception:
                            pass
                    memory.emit_event("slack", "response_sent", {
                        "text": (response or "")[:200], "thread_ts": thread_ts,
                    })
                    slack_text = md_to_slack(response)
                    await web_client.chat_postMessage(
                        channel=channel_id, text=slack_text, thread_ts=thread_ts)
                    memory.mark_message_processed(dedup_key, slack_ts, channel_id)
                    return

                # Parallel execution: spawn a new CLI session for this directive
                # Multiple CLIs can run concurrently within the same thread
                busy = cli_pool.busy_count(thread_ts)
                if busy > 0:
                    memory.emit_event("cli", "parallel_spawn", {
                        "thread_ts": thread_ts,
                        "existing_busy": busy,
                        "new_directive": directive_text[:200],
                    })

                # Update status: directive accepted
                if thinking_msg:
                    try:
                        await web_client.chat_update(
                            channel=channel_id, ts=thinking_msg,
                            text=f":rocket: *Directive accepted:* _{directive_text[:120]}_\n:brain: Gathering intelligence...")
                    except Exception:
                        pass

                # Build ML intelligence briefing + RAG context (parallel)
                ml_briefing = await asyncio.to_thread(_build_ml_briefing, directive_text)
                rag_context = await asyncio.to_thread(
                    build_rag_context, directive_text, 8000, {f"thread:{thread_ts}"},
                )

                # Sonnet planning step — creates structured execution plan
                if thinking_msg:
                    try:
                        await web_client.chat_update(
                            channel=channel_id, ts=thinking_msg,
                            text=f":brain: *Planning:* _{directive_text[:120]}_\n:memo: Sonnet creating execution plan...")
                    except Exception:
                        pass
                memory.emit_event("planner", "started", {
                    "thread_ts": thread_ts,
                    "directive": directive_text[:200],
                })

                execution_plan = await create_plan(
                    directive_text,
                    ml_briefing=ml_briefing,
                    rag_context=rag_context,
                    thread_history=thread_history,
                )

                if execution_plan:
                    memory.emit_event("planner", "completed", {
                        "thread_ts": thread_ts,
                        "plan_length": len(execution_plan),
                    })
                else:
                    memory.emit_event("planner", "skipped", {
                        "thread_ts": thread_ts,
                        "reason": "planner returned empty",
                    })

                # Update status: plan ready, agent starting
                briefing_summary = ""
                if ml_briefing and "similar directive" in ml_briefing.lower():
                    briefing_summary = "\n:mag: Found similar past work"
                plan_summary = "\n:memo: Execution plan ready" if execution_plan else ""
                if thinking_msg:
                    try:
                        await web_client.chat_update(
                            channel=channel_id, ts=thinking_msg,
                            text=f":robot_face: *Working:* _{directive_text[:120]}_"
                            f"{briefing_summary}{plan_summary}\n:hammer_and_wrench: Opus executing...")
                    except Exception:
                        pass

                # Build prompt with plan included
                cli_message = _build_threaded_prompt(
                    thread_history, directive_text,
                    ml_briefing=ml_briefing, rag_context=rag_context,
                    execution_plan=execution_plan,
                )
                cli_result = None
                try:
                    await safe_react("robot_face")
                    session = cli_pool.spawn(thread_ts, project_path)

                    memory.emit_event("cli", "started", {
                        "thread_ts": thread_ts,
                        "directive": directive_text[:200],
                        "pid": session.process.pid if session.process else None,
                    })

                    # Progress callback — updates Slack AND dashboard live
                    async def _progress(status: str) -> None:
                        if thinking_msg:
                            try:
                                await web_client.chat_update(
                                    channel=channel_id, ts=thinking_msg,
                                    text=f":robot_face: *Working:* _{directive_text[:80]}_\n{status}")
                            except Exception:
                                pass
                        memory.emit_event("cli", "progress", {
                            "thread_ts": thread_ts,
                            "status": status,
                        })

                    cli_result = await session.send(
                        cli_message, on_progress=_progress,
                    )

                    if cli_result.succeeded:
                        response = cli_result.output
                        if response and _is_garbage_response(response):
                            response = None
                        memory.emit_event("cli", "completed", {
                            "thread_ts": thread_ts,
                            "cost_usd": cli_result.cost_usd,
                            "elapsed_seconds": cli_result.elapsed_seconds,
                            "tools_used": cli_result.metadata.get("tools_used", 0),
                            "files_touched": cli_result.metadata.get("files_touched", []),
                            "output_len": len(response or ""),
                        })
                    else:
                        error_detail = cli_result.output or cli_result.error_detail or "Unknown error"
                        error_type = cli_result.error_type or "unknown"
                        memory.emit_event("cli", "error", {
                            "thread_ts": thread_ts,
                            "error_type": error_type,
                            "detail": error_detail[:300],
                        })
                        if error_type == "timeout":
                            response = f":warning: *Timed out* after 15 minutes.\n{error_detail[:1000]}"
                        elif error_type == "cancelled":
                            response = None
                        else:
                            response = f":warning: *Error:* {error_detail[:1000]}"

                except Exception as e:
                    logger.error("CLI error: %s", e)
                    response = f":warning: *Error:* {e}"

                if not response or response == "(No response from CLI)":
                    response = "I had trouble processing that request. Could you try rephrasing?"

                # === POST-EXECUTION SUMMARY REPORT ===
                # Send a concise execution report as a follow-up thread message
                if cli_result:
                    try:
                        report = _build_execution_report(cli_result, directive_text)
                        if report:
                            await web_client.chat_postMessage(
                                channel=channel_id,
                                text=report,
                                thread_ts=thread_ts,
                            )
                    except Exception as e:
                        logger.warning("Execution report error: %s", e)

            elif intake_result.tool_called == "generate_document":
                # Document generation — use existing generate_document flow
                tool_input = intake_result.tool_input or {}
                doc_type = tool_input.get("document_type", "pdf")
                doc_desc = tool_input.get("description", text)

                if thinking_msg:
                    await web_client.chat_update(
                        channel=channel_id, ts=thinking_msg,
                        text=f":page_facing_up: Generating {doc_type.upper()}...")

                doc_info = {"format": doc_type, "title": doc_desc[:100]}
                result = await generate_document(doc_desc, doc_info)
                if "error" in result:
                    response = f":warning: Document generation failed: {result['error']}"
                else:
                    await upload_file_to_slack(
                        web_client, channel_id, result["filepath"],
                        title=result.get("title", "Document"),
                        comment=f"Generated {result['format'].upper()}: *{result.get('title', 'Document')}*",
                        thread_ts=thread_ts)
                    try:
                        ingest_conversation(thread_ts, text, f"Generated {result['format']}: {result.get('title', 'Document')}")
                    except Exception:
                        pass
                    if thinking_msg:
                        try:
                            await web_client.chat_delete(channel=channel_id, ts=thinking_msg)
                        except Exception:
                            pass
                    memory.mark_message_processed(dedup_key, slack_ts, channel_id)
                    return  # Document uploaded, no text response needed

            elif intake_result.tool_called in ("mutate_org", "talk_to_agent"):
                # Org change or agent communication — dispatch handles it
                response = await dispatch(intake_result)

            else:
                # All other query tools (query_org, query_status, query_cost, query_kpi, query_ml)
                # Dispatch executes the query and returns formatted data
                raw_data = await dispatch(intake_result)

                # Send tool result back to Haiku for natural language formatting
                response = await format_with_tool_result(intake_result, raw_data, thread_history)

            memory.emit_event("slack", "response_sent", {
                "text": (response or "")[:200], "thread_ts": thread_ts,
            })
            if not response:
                response = ""
            slack_text = md_to_slack(response)

            # Use Block Kit for responses containing code
            has_code = '```' in response or '`' in response
            if has_code:
                blocks = format_code_output(response)
                if blocks:
                    await web_client.chat_postMessage(
                        channel=channel_id, text=slack_text[:3900],
                        blocks=blocks, thread_ts=thread_ts)
                else:
                    await web_client.chat_postMessage(
                        channel=channel_id, text=slack_text[:3900],
                        thread_ts=thread_ts)
            elif len(slack_text) > 3900:
                chunks = [slack_text[i:i+3900] for i in range(0, len(slack_text), 3900)]
                for chunk in chunks:
                    await web_client.chat_postMessage(
                        channel=channel_id, text=chunk,
                        thread_ts=thread_ts)
            else:
                await web_client.chat_postMessage(
                    channel=channel_id, text=slack_text,
                    thread_ts=thread_ts)

            # Remove the "Processing..." indicator
            if thinking_msg:
                try:
                    await web_client.chat_delete(channel=channel_id, ts=thinking_msg)
                except Exception:
                    pass

            # Ingest this exchange into RAG for future retrieval
            if response and not response.startswith("I had trouble"):
                try:
                    ingest_conversation(thread_ts, text, response)
                except Exception:
                    pass  # RAG ingestion is best-effort

            # Mark message as processed for idempotency
            memory.mark_message_processed(dedup_key, slack_ts, channel_id)

            logger.info("Responded in thread %s: %s...", thread_ts, slack_text[:100])

        except Exception as e:
            error_msg = f"{type(e).__name__}: {str(e)[:200]}"
            logger.error("Error: %s", error_msg, exc_info=True)
            # Remove thinking indicator on error
            if thinking_msg:
                try:
                    await web_client.chat_delete(channel=channel_id, ts=thinking_msg)
                except Exception:
                    pass
            try:
                await web_client.chat_postMessage(
                    channel=channel_id,
                    text=f":warning: Error processing your request:\n`{error_msg}`",
                    thread_ts=thread_ts,
                )
            except Exception:
                pass

    socket_client.socket_mode_request_listeners.append(handle_event)
    logger.info("Connecting...")
    await socket_client.connect()
    cli_pool.start_cleanup_loop()
    logger.info("Connected and listening. CLI session pool active.")

    while True:
        await asyncio.sleep(1)


if __name__ == "__main__":
    asyncio.run(start_slack_listener())
