"""
NEXUS Document Generator

Creates DOCX, PPTX, PDF, and image files from natural language requests.
Uses Gemini for content generation, then python-docx/python-pptx/reportlab/PIL
to create the actual files.
"""

import html
import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path

import google.genai as genai

logger = logging.getLogger("nexus.documents")

from src.config import get_key as _load_key  # consolidated key loading

# Security: Define allowed output directory
ALLOWED_OUTPUT_DIR = Path("~/.nexus/documents").expanduser()


def sanitize_filename(name: str) -> str:
    """Sanitize filename to prevent path traversal attacks.

    Removes path separators, null bytes, and restricts to alphanumeric characters
    plus dash and underscore. Limits length to 255 characters.
    """
    # Remove path traversal characters and null bytes
    safe = re.sub(r'[/\\.\0]', '', name)
    # Replace any non-alphanumeric (except dash/underscore) with underscore
    safe = re.sub(r'[^a-zA-Z0-9_-]', '_', safe)
    # Limit length to filesystem-safe maximum
    return safe[:255]


def sanitize_document_content(content: str) -> str:
    """Sanitize document content to prevent injection attacks.

    Escapes HTML/XML special characters and removes potentially dangerous
    script tags or macro content.
    """
    # Escape HTML/XML special characters
    content = html.escape(content)
    # Remove script tags (case-insensitive, including content)
    content = re.sub(r'<script.*?</script>', '', content, flags=re.IGNORECASE | re.DOTALL)
    return content


async def _classify_context_needs(message: str) -> list[str]:
    """Use LLM to determine which internal data sources are relevant to a request.

    Returns a list of source tags: "org", "costs", "architecture", "dependencies",
    "readme", "source_tree", "config".
    """
    try:
        from src.agents.base import allm_call
        from src.agents.org_chart import HAIKU

        prompt = f"""Which internal data sources would help generate this document?
Message: "{message}"

Available sources:
- "org" — org chart, team structure, agents, reporting hierarchy
- "costs" — model costs, token budgets, spending breakdown
- "architecture" — system design, APIs, services, auth, data flow, databases
- "dependencies" — Python libraries, requirements.txt
- "readme" — project overview, features, high-level description
- "source_tree" — directory/file structure of the codebase
- "config" — pyproject.toml, project metadata

Return a JSON array of relevant source tags. Include ALL sources that would produce
a complete, accurate document. When in doubt, include the source.

Example: ["architecture", "dependencies", "source_tree", "readme"]

Return ONLY the JSON array."""

        raw, _ = await allm_call(prompt, HAIKU, max_tokens=200)
        cleaned = raw.strip()
        if "```" in cleaned:
            cleaned = cleaned.split("```json")[-1].split("```")[0].strip() if "```json" in cleaned else cleaned.split("```")[1].strip()
        result = json.loads(cleaned)
        if isinstance(result, list):
            return [s for s in result if isinstance(s, str)]
    except Exception as e:
        logger.warning("Context classification failed: %s", e)

    # Fallback: include everything so we never produce an empty document
    return ["org", "architecture", "dependencies", "readme", "source_tree", "config"]


async def _gather_internal_context(message: str) -> str:
    """Gather internal NEXUS data relevant to a document request.

    Uses LLM classification to determine which data sources are needed,
    then loads them from disk and runtime state.
    """
    sources = await _classify_context_needs(message)
    context_parts: list[str] = []
    repo_root = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", ".."))

    if "org" in sources:
        try:
            from src.agents.org_chart import get_org_summary
            context_parts.append(
                "=== NEXUS INTERNAL DATA: Organization Chart ===\n"
                f"{get_org_summary()}\n"
                "=== END INTERNAL DATA ===\n"
            )
        except ImportError:
            # Fallback to file on disk
            org_path = os.path.join(repo_root, "ORG_CHART.md")
            if os.path.exists(org_path):
                with open(org_path) as f:
                    context_parts.append(
                        "=== NEXUS INTERNAL DATA: Organization Chart ===\n"
                        f"{f.read()}\n=== END INTERNAL DATA ===\n"
                    )

    if "costs" in sources:
        try:
            from src.agents.org_chart import MODEL_COSTS
            from src.agents.registry import registry
            model_counts: dict[str, int] = {}
            for agent in registry.get_active_agents():
                m = agent.model
                model_counts[str(m)] = model_counts.get(str(m), 0) + 1
            summary = "Model Distribution:\n"
            for model, count in sorted(model_counts.items()):
                cost = MODEL_COSTS.get(model, {})
                summary += f"  {model}: {count} agents (input: ${cost.get('input', '?')}/M, output: ${cost.get('output', '?')}/M)\n"
            context_parts.append(
                "=== NEXUS INTERNAL DATA: Model Costs ===\n"
                f"{summary}=== END INTERNAL DATA ===\n"
            )
        except ImportError:
            pass

    if "architecture" in sources:
        arch_path = os.path.join(repo_root, "docs", "ARCHITECTURE.md")
        if os.path.exists(arch_path):
            with open(arch_path) as f:
                context_parts.append(
                    "=== NEXUS INTERNAL DATA: Architecture Documentation ===\n"
                    f"{f.read()}\n=== END INTERNAL DATA ===\n"
                )

    if "dependencies" in sources:
        req_path = os.path.join(repo_root, "requirements.txt")
        if os.path.exists(req_path):
            with open(req_path) as f:
                context_parts.append(
                    "=== NEXUS INTERNAL DATA: Python Dependencies ===\n"
                    f"{f.read()}\n=== END INTERNAL DATA ===\n"
                )

    if "readme" in sources:
        readme_path = os.path.join(repo_root, "README.md")
        if os.path.exists(readme_path):
            with open(readme_path) as f:
                context_parts.append(
                    "=== NEXUS INTERNAL DATA: README ===\n"
                    f"{f.read()}\n=== END INTERNAL DATA ===\n"
                )

    if "config" in sources:
        pyproject_path = os.path.join(repo_root, "pyproject.toml")
        if os.path.exists(pyproject_path):
            with open(pyproject_path) as f:
                context_parts.append(
                    "=== NEXUS INTERNAL DATA: Project Config ===\n"
                    f"{f.read()}\n=== END INTERNAL DATA ===\n"
                )

    if "source_tree" in sources:
        src_dir = os.path.join(repo_root, "src")
        if os.path.isdir(src_dir):
            tree_lines: list[str] = []
            for root, dirs, files in os.walk(src_dir):
                dirs[:] = [d for d in sorted(dirs) if d != "__pycache__"]
                depth = root.replace(src_dir, "").count(os.sep)
                indent = "  " * depth
                tree_lines.append(f"{indent}{os.path.basename(root)}/")
                for f_name in sorted(files):
                    if f_name.endswith(".py"):
                        tree_lines.append(f"{indent}  {f_name}")
            context_parts.append(
                "=== NEXUS INTERNAL DATA: Source Tree ===\n"
                + "\n".join(tree_lines) + "\n"
                "=== END INTERNAL DATA ===\n"
            )

    return "\n".join(context_parts)


def _needs_web_enrichment(message: str, has_internal_context: bool) -> str | None:
    """Determine if a request needs web search and return a search query if so.

    Returns a search query string, or None if web search isn't needed.
    """
    msg_lower = message.lower()

    # Explicit web indicators — always search
    explicit = ["compare to", "compared to", "vs ", "versus", "industry",
                "benchmark", "best practice", "how does", "market",
                "competitor", "trend", "research", "statistics", "data on",
                "according to", "report on", "analysis of"]
    if any(kw in msg_lower for kw in explicit):
        # Build a focused search query from the request
        # Strip NEXUS-internal references to get a cleaner web query
        clean = msg_lower
        for strip in ["nexus", "our team", "our org", "my team", "our company",
                       "pdf", "docx", "pptx", "slides", "document", "report",
                       "create", "generate", "make", "send", "give me", "show me"]:
            clean = clean.replace(strip, "")
        return clean.strip()[:200] or None

    # If we have NO internal context and the request is about something specific,
    # it's probably about external info — search for it
    if not has_internal_context:
        # Topics that are almost certainly external
        external_signals = ["how to", "what is", "who is", "guide", "tutorial",
                           "example", "template", "standard", "framework",
                           "technology", "platform", "tool", "service"]
        if any(kw in msg_lower for kw in external_signals):
            return msg_lower[:200]

    return None


async def _gather_web_context(query: str) -> str:
    """Search the web and return results formatted as context for document generation."""
    try:
        from src.tools.web_search import format_results_for_context
        from src.tools.web_search import search as web_search
        results = await web_search(query, num_results=5)
        if results:
            formatted = format_results_for_context(results)
            return (
                "=== WEB RESEARCH RESULTS ===\n"
                f"{formatted}\n"
                "=== END WEB RESEARCH ===\n"
            )
    except Exception as e:
        logger.warning("Web search failed (non-fatal): %s", e)
    return ""


def _ask_gemini(prompt: str, system: str = "") -> str:
    """Generate content using Gemini API, with Claude fallback on rate limit."""
    # Try Gemini first
    api_key = _load_key("GOOGLE_AI_API_KEY")
    if api_key:
        try:
            client = genai.Client(api_key=api_key)
            full_prompt = f"{system}\n\n{prompt}" if system else prompt
            response = client.models.generate_content(
                model="gemini-2.5-pro",
                contents=full_prompt,
            )
            return str(response.text)
        except Exception as e:
            logger.warning("Gemini failed (%s), falling back to Claude", e)

    # Claude fallback
    anthropic_key = _load_key("ANTHROPIC_API_KEY")
    if not anthropic_key:
        raise ValueError("No AI API keys available (tried Gemini and Claude)")

    import anthropic
    claude_client = anthropic.Anthropic(api_key=anthropic_key)
    messages = [{"role": "user", "content": prompt}]
    claude_response = claude_client.messages.create(
        model="claude-sonnet-4-6-20250929",
        max_tokens=4096,
        system=system or "You generate document content as requested.",
        messages=messages,  # type: ignore[arg-type]
    )
    return str(claude_response.content[0].text)  # type: ignore[union-attr]


def _parse_json(content: str) -> dict:
    """Parse JSON content, stripping markdown fences if present."""
    content = content.strip()
    if content.startswith("```"):
        lines = content.split("\n")
        # Remove first line (```json or ```)
        content = "\n".join(lines[1:])
    if content.endswith("```"):
        content = content.rsplit("```", 1)[0]
    content = content.strip()
    return json.loads(content)  # type: ignore[no-any-return]


# ============================================
# DOCX
# ============================================

def _parse_markdown_to_docx_elements(markdown_text: str) -> list[dict]:
    """Parse markdown text into structured elements for DOCX rendering.

    Handles headings (#), bullets (- or *), numbered lists (1.), code blocks (```),
    tables (| ... |), blockquotes (>), and plain paragraphs.
    """
    import re

    elements: list[dict] = []
    lines = markdown_text.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            i += 1
            continue

        # Code block (```)
        if stripped.startswith("```"):
            language = stripped[3:].strip()
            code_lines: list[str] = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            elements.append({"type": "code", "code": "\n".join(code_lines), "language": language})
            continue

        # Headings (# through ####)
        heading_match = re.match(r'^(#{1,4})\s+(.+)$', stripped)
        if heading_match:
            level = len(heading_match.group(1))
            elements.append({"type": "heading", "text": heading_match.group(2).strip(), "level": level})
            i += 1
            continue

        # Table (| Header | Header |)
        if stripped.startswith("|") and "|" in stripped[1:]:
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            # Parse table
            headers: list[str] = []
            rows: list[list[str]] = []
            for _ti, tline in enumerate(table_lines):
                cells = [c.strip() for c in tline.strip("|").split("|")]
                # Skip separator rows (|---|---|)
                if all(re.match(r'^[-:]+$', c) for c in cells):
                    continue
                if not headers:
                    headers = cells
                else:
                    rows.append(cells)
            if headers:
                elements.append({"type": "table", "headers": headers, "rows": rows})
            continue

        # Blockquote (>)
        if stripped.startswith(">"):
            quote_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(lines[i].strip().lstrip(">").strip())
                i += 1
            elements.append({"type": "quote", "text": " ".join(quote_lines)})
            continue

        # Bullet list (- or *)
        if re.match(r'^[\-\*]\s+', stripped):
            items: list[str] = []
            while i < len(lines) and re.match(r'^\s*[\-\*]\s+', lines[i]):
                items.append(re.sub(r'^\s*[\-\*]\s+', '', lines[i]).strip())
                i += 1
            elements.append({"type": "bullets", "items": items})
            continue

        # Numbered list (1. 2. etc.)
        if re.match(r'^\d+[\.\)]\s+', stripped):
            items = []
            while i < len(lines) and re.match(r'^\s*\d+[\.\)]\s+', lines[i]):
                items.append(re.sub(r'^\s*\d+[\.\)]\s+', '', lines[i]).strip())
                i += 1
            elements.append({"type": "numbered_list", "items": items})
            continue

        # Bold-prefixed lines that act as inline sub-headings (e.g., "**Key Point:** ...")
        if stripped.startswith("**") and "**" in stripped[2:]:
            elements.append({"type": "paragraph", "text": stripped})
            i += 1
            continue

        # Plain paragraph — collect consecutive non-empty, non-special lines
        para_lines: list[str] = []
        while i < len(lines):
            cur = lines[i].strip()
            if not cur:
                i += 1
                break
            # Stop if next line is a special element
            if (cur.startswith("#") or cur.startswith("```") or cur.startswith("|")
                    or cur.startswith(">") or re.match(r'^[\-\*]\s+', cur)
                    or re.match(r'^\d+[\.\)]\s+', cur)):
                break
            para_lines.append(cur)
            i += 1
        if para_lines:
            elements.append({"type": "paragraph", "text": " ".join(para_lines)})
        continue

    return elements


def create_docx(title: str, request: str, output_dir: str | None = None) -> str:
    """Generate a Word document based on the request.

    Uses a markdown-based generation approach: the LLM produces well-formatted
    markdown prose, which is then parsed into DOCX elements. This avoids the
    fragile JSON serialization that caused raw JSON to appear in documents.
    """
    import re

    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor

    content = _ask_gemini(
        f"Write the full content for a professional document.\n\n"
        f"Title: {title}\n"
        f"Request: {request}\n\n"
        f"FORMAT RULES:\n"
        f"- Write in well-structured markdown with proper headings (# for top-level, ## for sub-sections, ### for details).\n"
        f"- Use bullet lists (- item), numbered lists (1. item), tables (| col | col |), and code blocks (```) where appropriate.\n"
        f"- Write substantive prose paragraphs — not bullet-only content. Each section should have explanatory text.\n"
        f"- Do NOT wrap output in JSON. Do NOT use code fences around the entire document.\n"
        f"- Do NOT include ```markdown or ```json wrappers.\n"
        f"- Output ONLY the markdown document content, starting with the first # heading.\n",
        system="You are a professional technical writer. Write detailed, well-structured markdown documents. "
               "Never output JSON. Never wrap your response in code fences. Write natural prose with markdown formatting. "
               "If the request includes INTERNAL DATA sections, use that real data as the basis for the document content. "
               "Do NOT invent or hallucinate information when real data is provided."
    )

    # Strip any accidental markdown code fence wrapping
    content = content.strip()
    if content.startswith("```markdown"):
        content = content[len("```markdown"):].strip()
    if content.startswith("```md"):
        content = content[len("```md"):].strip()
    if content.startswith("```"):
        content = content[3:].strip()
    if content.endswith("```"):
        content = content[:-3].strip()

    # Parse the markdown into structured elements
    elements = _parse_markdown_to_docx_elements(content)

    doc = Document()

    # Set default fonts
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # Cover page
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(title)
    title_run.font.name = 'Cambria'
    title_run.font.size = Pt(28)
    title_run.font.bold = True

    doc.add_paragraph("")
    doc.add_paragraph("")

    generated_para = doc.add_paragraph()
    generated_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    generated_run = generated_para.add_run(f"Generated by NEXUS\n{datetime.now().strftime('%B %d, %Y')}")
    generated_run.font.size = Pt(10)
    generated_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # Page break after cover
    doc.add_page_break()

    # Table of contents placeholder
    toc_para = doc.add_paragraph("Table of Contents")
    toc_para.style = doc.styles['Heading 1']
    toc_field = doc.add_paragraph()
    toc_field.add_run("(Table of contents will be generated when opened in Word)")
    toc_field.paragraph_format.space_after = Pt(12)
    doc.add_page_break()

    # Section counter for numbering
    section_counters = {1: 0, 2: 0, 3: 0, 4: 0}

    def _strip_markdown_inline(text: str) -> list[tuple[str, bool, bool]]:
        """Split text into runs with bold/italic flags for inline markdown.

        Returns list of (text, is_bold, is_italic) tuples.
        """
        runs: list[tuple[str, bool, bool]] = []
        # Pattern: **bold**, *italic*, ***bold+italic***
        pattern = re.compile(r'(\*{1,3})(.*?)\1')
        pos = 0
        for m in pattern.finditer(text):
            # Add preceding plain text
            if m.start() > pos:
                runs.append((text[pos:m.start()], False, False))
            stars = len(m.group(1))
            inner = m.group(2)
            if stars == 3:
                runs.append((inner, True, True))
            elif stars == 2:
                runs.append((inner, True, False))
            else:
                runs.append((inner, False, True))
            pos = m.end()
        # Remaining text
        if pos < len(text):
            runs.append((text[pos:], False, False))
        return runs if runs else [(text, False, False)]

    def _add_rich_paragraph(doc_ref, text: str, style_name: str | None = None):
        """Add a paragraph with inline bold/italic markdown rendered."""
        para = doc_ref.add_paragraph(style=style_name)
        para.paragraph_format.space_after = Pt(6)
        for run_text, is_bold, is_italic in _strip_markdown_inline(text):
            run = para.add_run(run_text)
            if is_bold:
                run.font.bold = True
            if is_italic:
                run.font.italic = True
        return para

    # Render elements
    for elem in elements:
        elem_type = elem.get("type")

        if elem_type == "heading":
            level = min(elem.get("level", 1), 3)
            text = elem.get("text", "")

            # Update section numbering
            section_counters[level] += 1
            for reset_level in range(level + 1, 5):
                section_counters[reset_level] = 0

            # Build section number
            if level == 1:
                section_num = f"{section_counters[1]}. "
            elif level == 2:
                section_num = f"{section_counters[1]}.{section_counters[2]} "
            else:
                section_num = f"{section_counters[1]}.{section_counters[2]}.{section_counters[3]} "

            heading = doc.add_heading(section_num + text, level=level)
            heading.paragraph_format.space_before = Pt(12)
            heading.paragraph_format.space_after = Pt(6)
            if heading.runs:
                heading.runs[0].font.name = 'Cambria'

        elif elem_type == "paragraph":
            _add_rich_paragraph(doc, elem.get("text", ""))

        elif elem_type == "bullets":
            for bullet_text in elem.get("items", []):
                para = doc.add_paragraph(style='List Bullet')
                para.paragraph_format.space_after = Pt(3)
                for run_text, is_bold, is_italic in _strip_markdown_inline(bullet_text):
                    run = para.add_run(run_text)
                    if is_bold:
                        run.font.bold = True
                    if is_italic:
                        run.font.italic = True
            doc.add_paragraph("")

        elif elem_type == "numbered_list":
            for num_text in elem.get("items", []):
                para = doc.add_paragraph(style='List Number')
                para.paragraph_format.space_after = Pt(3)
                for run_text, is_bold, is_italic in _strip_markdown_inline(num_text):
                    run = para.add_run(run_text)
                    if is_bold:
                        run.font.bold = True
                    if is_italic:
                        run.font.italic = True
            doc.add_paragraph("")

        elif elem_type == "table":
            headers = elem.get("headers", [])
            rows = elem.get("rows", [])

            if headers:
                num_rows = 1 + len(rows)
                table = doc.add_table(rows=num_rows, cols=len(headers))
                table.style = 'Light Grid Accent 1'

                # Header row
                header_cells = table.rows[0].cells
                for ci, header in enumerate(headers):
                    header_cells[ci].text = header
                    for paragraph in header_cells[ci].paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

                # Data rows
                for row_idx, row_data in enumerate(rows):
                    row_cells = table.rows[row_idx + 1].cells
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < len(row_cells):
                            row_cells[col_idx].text = str(cell_data)

                doc.add_paragraph("")

        elif elem_type == "code":
            code_para = doc.add_paragraph()
            code_run = code_para.add_run(elem.get("code", ""))
            code_run.font.name = 'Courier New'
            code_run.font.size = Pt(9)

            shading_elm = OxmlElement('w:shd')
            shading_elm.set(qn('w:fill'), 'F0F0F5')
            code_para._element.get_or_add_pPr().append(shading_elm)

            code_para.paragraph_format.space_after = Pt(6)
            code_para.paragraph_format.left_indent = Pt(20)

        elif elem_type == "quote":
            quote_para = doc.add_paragraph()
            quote_para.paragraph_format.left_indent = Pt(40)
            quote_para.paragraph_format.space_after = Pt(6)

            quote_run = quote_para.add_run(f'"{elem.get("text", "")}"')
            quote_run.font.italic = True

    # Add header
    doc_section = doc.sections[0]
    header = doc_section.header
    header_para = header.paragraphs[0]
    header_run = header_para.add_run(title)
    header_run.font.size = Pt(9)
    header_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # Add footer with page number
    footer = doc_section.footer
    footer_para = footer.paragraphs[0]
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_para.add_run("Page ")

    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')

    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = "PAGE"

    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'end')

    footer_para._element.append(fldChar1)
    footer_para._element.append(instrText)
    footer_para._element.append(fldChar2)

    # Security: Sanitize output directory and filename
    if output_dir:
        output_path = Path(output_dir).resolve()
        # Verify output path is within allowed directory
        try:
            output_path.relative_to(ALLOWED_OUTPUT_DIR)
        except ValueError:
            raise ValueError(f"Output directory must be within {ALLOWED_OUTPUT_DIR}") from None  # noqa: B904
    else:
        output_path = ALLOWED_OUTPUT_DIR
        output_path.mkdir(parents=True, exist_ok=True)

    safe_title = sanitize_filename(title[:50])
    filepath = output_path / f"{safe_title}.docx"
    doc.save(str(filepath))
    return str(filepath)


# ============================================
# PPTX
# ============================================

def create_pptx(title: str, request: str, output_dir: str | None = None) -> str:
    """Generate a PowerPoint presentation based on the request."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    content = _ask_gemini(
        f"Generate content for a professional slide deck.\n\n"
        f"Title: {title}\n"
        f"Request: {request}\n\n"
        f"Return as JSON with this structure:\n"
        f'{{"title": "...", "subtitle": "...", "slides": ['
        f'{{"type": "content", "title": "...", "bullets": ["...", "..."], "notes": "..."}}, '
        f'{{"type": "two_column", "title": "...", "left": ["...", "..."], "right": ["...", "..."], "notes": "..."}}, '
        f'{{"type": "comparison", "title": "...", "left_header": "...", "left": ["..."], "right_header": "...", "right": ["..."], "notes": "..."}}, '
        f'{{"type": "quote", "quote": "...", "attribution": "...", "notes": "..."}}, '
        f'{{"type": "closing", "message": "Thank you", "submessage": "Questions?", "notes": "..."}}'
        f']}}\n\n'
        f"Generate 6-12 slides. Use a mix of slide types. Only return the JSON, nothing else.",
        system="You generate presentation content as structured JSON. No markdown, no code fences, just JSON. "
               "If the request includes INTERNAL DATA sections, use that real data as the basis for the slides. "
               "Do NOT invent or hallucinate information when real data is provided."
    )

    try:
        data = _parse_json(content)
    except json.JSONDecodeError:
        data = {
            "title": title,
            "subtitle": request,
            "slides": [{"type": "content", "title": "Content", "bullets": [request], "notes": ""}],
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    bg_color = RGBColor(0x1A, 0x1A, 0x2E)
    accent_color = RGBColor(0x00, 0xD2, 0xFF)
    text_color = RGBColor(0xFF, 0xFF, 0xFF)
    subtitle_color = RGBColor(0xAA, 0xAA, 0xCC)

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_number(slide, slide_num):
        """Add slide number to bottom right."""
        txBox = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = Pt(10)
        p.font.color.rgb = subtitle_color
        p.alignment = PP_ALIGN.RIGHT

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, bg_color)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", title)
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    if data.get("subtitle"):
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = data["subtitle"]
        p2.font.size = Pt(20)
        p2.font.color.rgb = subtitle_color
        p2.alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(1, Inches(4), Inches(4), Inches(5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # Content slides
    slide_num = 1
    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type", "content")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, bg_color)

        # Add slide number
        add_slide_number(slide, slide_num)
        slide_num += 1

        # Title (common to most slide types)
        if slide_type in ["content", "two_column", "comparison"]:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "")
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = accent_color

            # Accent underline
            line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(2), Pt(3))
            line.fill.solid()
            line.fill.fore_color.rgb = accent_color
            line.line.fill.background()

        if slide_type == "content":
            # Bullets
            bullets = slide_data.get("bullets", [])
            if bullets:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    # Handle sub-bullets (items starting with "  - " or similar)
                    if isinstance(bullet, str) and bullet.strip().startswith("-"):
                        p.text = f"  ▹  {bullet.lstrip('- ').strip()}"
                        p.level = 1
                        p.font.size = Pt(16)
                    else:
                        p.text = f"▸  {bullet}"
                        p.level = 0
                        p.font.size = Pt(20)

                    p.font.color.rgb = text_color
                    p.space_after = Pt(12)

        elif slide_type == "two_column":
            # Two columns side by side
            left_items = slide_data.get("left", [])
            right_items = slide_data.get("right", [])

            # Left column
            txBox_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
            tf_left = txBox_left.text_frame
            tf_left.word_wrap = True

            for i, item in enumerate(left_items):
                if i == 0:
                    p = tf_left.paragraphs[0]
                else:
                    p = tf_left.add_paragraph()
                p.text = f"▸  {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_after = Pt(10)

            # Right column
            txBox_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(5))
            tf_right = txBox_right.text_frame
            tf_right.word_wrap = True

            for i, item in enumerate(right_items):
                if i == 0:
                    p = tf_right.paragraphs[0]
                else:
                    p = tf_right.add_paragraph()
                p.text = f"▸  {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_after = Pt(10)

        elif slide_type == "comparison":
            # Comparison with headers
            left_header = slide_data.get("left_header", "Option A")
            right_header = slide_data.get("right_header", "Option B")
            left_items = slide_data.get("left", [])
            right_items = slide_data.get("right", [])

            # Left header
            txBox_lh = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5))
            tf_lh = txBox_lh.text_frame
            p_lh = tf_lh.paragraphs[0]
            p_lh.text = left_header
            p_lh.font.size = Pt(24)
            p_lh.font.bold = True
            p_lh.font.color.rgb = accent_color

            # Left items
            txBox_l = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.5))
            tf_l = txBox_l.text_frame
            tf_l.word_wrap = True

            for i, item in enumerate(left_items):
                if i == 0:
                    p = tf_l.paragraphs[0]
                else:
                    p = tf_l.add_paragraph()
                p.text = f"▸  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
                p.space_after = Pt(8)

            # Right header
            txBox_rh = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.5))
            tf_rh = txBox_rh.text_frame
            p_rh = tf_rh.paragraphs[0]
            p_rh.text = right_header
            p_rh.font.size = Pt(24)
            p_rh.font.bold = True
            p_rh.font.color.rgb = accent_color

            # Right items
            txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(5.5), Inches(4.5))
            tf_r = txBox_r.text_frame
            tf_r.word_wrap = True

            for i, item in enumerate(right_items):
                if i == 0:
                    p = tf_r.paragraphs[0]
                else:
                    p = tf_r.add_paragraph()
                p.text = f"▸  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
                p.space_after = Pt(8)

        elif slide_type == "quote":
            # Large centered quote
            quote_text = slide_data.get("quote", "")
            attribution = slide_data.get("attribution", "")

            txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{quote_text}"'
            p.font.size = Pt(32)
            p.font.italic = True
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

            if attribution:
                txBox_attr = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(0.8))
                tf_attr = txBox_attr.text_frame
                p_attr = tf_attr.paragraphs[0]
                p_attr.text = f"— {attribution}"
                p_attr.font.size = Pt(20)
                p_attr.font.color.rgb = subtitle_color
                p_attr.alignment = PP_ALIGN.CENTER

        elif slide_type == "closing":
            # Thank you / closing slide
            message = slide_data.get("message", "Thank you")
            submessage = slide_data.get("submessage", "Questions?")

            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = message
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = submessage
            p2.font.size = Pt(24)
            p2.font.color.rgb = accent_color
            p2.alignment = PP_ALIGN.CENTER

        # Speaker notes
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    # Security: Sanitize output directory and filename
    if output_dir:
        output_path = Path(output_dir).resolve()
        # Verify output path is within allowed directory
        try:
            output_path.relative_to(ALLOWED_OUTPUT_DIR)
        except ValueError:
            raise ValueError(f"Output directory must be within {ALLOWED_OUTPUT_DIR}") from None  # noqa: B904
    else:
        output_path = ALLOWED_OUTPUT_DIR
        output_path.mkdir(parents=True, exist_ok=True)

    safe_title = sanitize_filename(title[:50])
    filepath = output_path / f"{safe_title}.pptx"
    prs.save(str(filepath))
    return str(filepath)


# ============================================
# PDF
# ============================================

def create_pdf(title: str, request: str, output_dir: str | None = None) -> str:
    """Generate a PDF document based on the request."""
    from reportlab.lib.colors import HexColor
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        ListFlowable,
        ListItem,
        PageBreak,
        Paragraph,
        Preformatted,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    content = _ask_gemini(
        f"Generate the full content for a professional document.\n\n"
        f"Title: {title}\n"
        f"Request: {request}\n\n"
        f"Return as JSON with this structure:\n"
        f'{{"title": "...", "subtitle": "...", "sections": ['
        f'{{"heading": "...", "level": 1, "content": ['
        f'{{"type": "paragraph", "text": "..."}}, '
        f'{{"type": "bullets", "items": ["...", "..."]}}, '
        f'{{"type": "numbered_list", "items": ["...", "..."]}}, '
        f'{{"type": "table", "headers": ["..."], "rows": [["...", "..."]]}} '
        f'{{"type": "code", "language": "python", "code": "..."}}, '
        f'{{"type": "quote", "text": "...", "attribution": "..."}}'
        f']}}, ...]}}\n\n'
        f"Only return the JSON, nothing else.",
        system="You generate document content as structured JSON. No markdown, no code fences, just JSON. "
               "If the request includes INTERNAL DATA sections, use that real data as the basis for the document content. "
               "Do NOT invent or hallucinate information when real data is provided."
    )

    try:
        data = _parse_json(content)
    except json.JSONDecodeError:
        data = {
            "title": title,
            "subtitle": "",
            "sections": [{"heading": "Content", "level": 1, "content": [{"type": "paragraph", "text": content}]}],
        }

    # Security: Sanitize output directory and filename
    if output_dir:
        output_path = Path(output_dir).resolve()
        # Verify output path is within allowed directory
        try:
            output_path.relative_to(ALLOWED_OUTPUT_DIR)
        except ValueError:
            raise ValueError(f"Output directory must be within {ALLOWED_OUTPUT_DIR}") from None  # noqa: B904
    else:
        output_path = ALLOWED_OUTPUT_DIR
        output_path.mkdir(parents=True, exist_ok=True)

    safe_title = sanitize_filename(title[:50])
    filepath = str(output_path / f"{safe_title}.pdf")

    # Custom page template with header/footer
    def add_page_elements(canvas, doc):
        canvas.saveState()

        # Header
        canvas.setFont('Helvetica', 9)
        canvas.setFillColor(HexColor("#999999"))
        canvas.drawString(inch, letter[1] - 0.5*inch, data.get("title", title))

        # Page number
        page_num = canvas.getPageNumber()
        canvas.drawRightString(letter[0] - inch, 0.5*inch, f"Page {page_num}")

        canvas.restoreState()

    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        rightMargin=inch,
        leftMargin=inch,
        topMargin=inch,
        bottomMargin=inch,
    )

    styles = getSampleStyleSheet()

    # Custom styles
    styles.add(ParagraphStyle(
        name="DocTitle",
        parent=styles["Title"],
        fontSize=24,
        spaceAfter=6,
        textColor=HexColor("#1a1a2e"),
        alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="DocSubtitle",
        parent=styles["Normal"],
        fontSize=14,
        spaceAfter=20,
        textColor=HexColor("#666688"),
        alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="SectionHead1",
        parent=styles["Heading1"],
        fontSize=18,
        spaceBefore=12,
        spaceAfter=6,
        textColor=HexColor("#00d2ff"),
    ))
    styles.add(ParagraphStyle(
        name="SectionHead2",
        parent=styles["Heading2"],
        fontSize=14,
        spaceBefore=10,
        spaceAfter=5,
        textColor=HexColor("#00d2ff"),
    ))
    styles.add(ParagraphStyle(
        name="BodyText2",
        parent=styles["Normal"],
        fontSize=11,
        spaceAfter=6,
        leading=14,
    ))
    styles.add(ParagraphStyle(
        name="CodeBlock",
        parent=styles["Code"],
        fontName="Courier",
        fontSize=9,
        leftIndent=20,
        spaceAfter=6,
        backColor=HexColor("#f0f0f5"),
    ))
    styles.add(ParagraphStyle(
        name="QuoteBlock",
        parent=styles["Normal"],
        fontSize=11,
        leftIndent=40,
        spaceAfter=6,
        fontName="Helvetica-Oblique",
    ))

    story = []

    # Cover page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph(data.get("title", title), styles["DocTitle"]))

    if data.get("subtitle"):
        story.append(Paragraph(data["subtitle"], styles["DocSubtitle"]))

    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(f"Generated by NEXUS<br/>{datetime.now().strftime('%B %d, %Y')}",
                          ParagraphStyle(name="Generated", parent=styles["Normal"],
                                       fontSize=10, textColor=HexColor("#999999"), alignment=TA_CENTER)))
    story.append(PageBreak())

    # Table of contents placeholder
    story.append(Paragraph("Table of Contents", styles["SectionHead1"]))

    toc_items = []
    section_num = 0
    for section in data.get("sections", []):
        if section.get("heading"):
            section_num += 1
            toc_items.append([f"{section_num}.", section["heading"], ""])

    if toc_items:
        toc_table = Table(toc_items, colWidths=[0.5*inch, 5*inch, 0.5*inch])
        toc_table.setStyle(TableStyle([
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('TEXTCOLOR', (0, 0), (-1, -1), HexColor("#333333")),
            ('LEFTPADDING', (0, 0), (-1, -1), 0),
            ('RIGHTPADDING', (0, 0), (-1, -1), 0),
        ]))
        story.append(toc_table)

    story.append(PageBreak())

    # Section counter for numbering
    section_num = 0

    # Sections
    for section in data.get("sections", []):
        level = section.get("level", 1)

        if section.get("heading"):
            section_num += 1
            section_text = f"{section_num}. {section['heading']}"

            if level == 1:
                story.append(Paragraph(section_text, styles["SectionHead1"]))
            else:
                story.append(Paragraph(section_text, styles["SectionHead2"]))

        # Render content items
        for item in section.get("content", []):
            item_type = item.get("type")

            if item_type == "paragraph":
                text = item.get("text", "")
                # Escape XML-sensitive characters
                safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe, styles["BodyText2"]))

            elif item_type == "bullets":
                bullet_items = []
                for bullet_text in item.get("items", []):
                    safe = bullet_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    bullet_items.append(ListItem(Paragraph(safe, styles["BodyText2"]),
                                                leftIndent=20, bulletColor=HexColor("#00d2ff")))
                story.append(ListFlowable(bullet_items, bulletType='bullet'))
                story.append(Spacer(1, 6))

            elif item_type == "numbered_list":
                num_items = []
                for num_text in item.get("items", []):
                    safe = num_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    num_items.append(ListItem(Paragraph(safe, styles["BodyText2"]), leftIndent=20))
                story.append(ListFlowable(num_items, bulletType='1'))
                story.append(Spacer(1, 6))

            elif item_type == "table":
                headers = item.get("headers", [])
                rows = item.get("rows", [])

                if headers and rows:
                    table_data = [headers] + rows
                    t = Table(table_data)
                    t.setStyle(TableStyle([
                        # Header row
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor("#00d2ff")),
                        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#ffffff")),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 11),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),

                        # Data rows
                        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 1), (-1, -1), 10),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [HexColor("#ffffff"), HexColor("#f5f5f5")]),

                        # Grid
                        ('GRID', (0, 0), (-1, -1), 0.5, HexColor("#cccccc")),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 12))

            elif item_type == "code":
                code_text = item.get("code", "")
                story.append(Preformatted(code_text, styles["CodeBlock"]))

            elif item_type == "quote":
                quote_text = item.get("text", "")
                attribution = item.get("attribution", "")

                # Create quote with left border effect using table
                quote_content = f'"{quote_text}"'
                if attribution:
                    quote_content += f"\n— {attribution}"

                safe = quote_content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

                quote_table = Table([[Paragraph(safe, styles["QuoteBlock"])]], colWidths=[5*inch])
                quote_table.setStyle(TableStyle([
                    ('LEFTPADDING', (0, 0), (0, 0), 10),
                    ('RIGHTPADDING', (0, 0), (0, 0), 10),
                    ('TOPPADDING', (0, 0), (0, 0), 6),
                    ('BOTTOMPADDING', (0, 0), (0, 0), 6),
                    ('LINEAFTER', (0, 0), (0, 0), 3, HexColor("#00d2ff")),
                ]))
                story.append(quote_table)
                story.append(Spacer(1, 6))

    doc.build(story, onFirstPage=add_page_elements, onLaterPages=add_page_elements)
    return filepath


# ============================================
# IMAGE
# ============================================

def _extract_image_prompt(raw_description: str) -> str:
    """Extract a clean, focused image prompt from a raw request.

    Strips internal context blocks, web research sections, and meta-instructions
    so the image model gets a clear visual description.
    """
    # Strip NEXUS internal data sections
    import re
    clean = re.sub(r'===\s*(NEXUS INTERNAL DATA|WEB RESEARCH|END INTERNAL|END WEB).*?===\s*',
                   '', raw_description, flags=re.DOTALL)
    # Strip enrichment instructions
    clean = re.sub(r'IMPORTANT:.*?(?=\n\n|\Z)', '', clean, flags=re.DOTALL)
    clean = re.sub(r'Prefer internal data.*?(?=\n\n|\Z)', '', clean, flags=re.DOTALL)
    clean = re.sub(r'Do NOT make up.*?(?=\n\n|\Z)', '', clean, flags=re.DOTALL)
    clean = clean.strip()

    # If still too long, ask LLM to distill into an image prompt
    if len(clean) > 500:
        try:
            import asyncio

            from src.agents.base import allm_call
            from src.agents.org_chart import HAIKU
            loop = asyncio.get_event_loop()
            summary, _ = loop.run_until_complete(allm_call(
                f"Distill this into a concise image generation prompt (1-3 sentences). "
                f"Focus on what the image should LOOK like visually:\n\n{clean[:2000]}",
                HAIKU, max_tokens=200))
            return summary.strip()
        except Exception:
            pass

    return clean[:500] if clean else raw_description[:500]


def _try_gemini_image_generation(description: str, output_path: str) -> bool:
    """Generate an image using Gemini's native image generation.

    Tries models in order: gemini-2.5-flash-image (production),
    gemini-3-pro-image-preview (pro quality).
    Returns True if successful, False to fall back to PIL.
    """
    api_key = _load_key("GOOGLE_AI_API_KEY")
    if not api_key:
        return False

    # Clean the prompt before sending to image models
    clean_prompt = _extract_image_prompt(description)
    logger.info("Image prompt: %s...", clean_prompt[:100])

    models = ["gemini-2.5-flash-image", "gemini-3-pro-image-preview"]

    for model_name in models:
        try:
            from google.genai import types

            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model=model_name,
                contents=clean_prompt,
                config=types.GenerateContentConfig(
                    response_modalities=["TEXT", "IMAGE"],
                ),
            )

            for part in response.parts:  # type: ignore[union-attr]
                if part.inline_data is not None:
                    image = part.as_image()
                    image.save(output_path)  # type: ignore[union-attr]
                    logger.info("Image generated with %s", model_name)
                    return True

            logger.warning("%s returned no image data", model_name)

        except Exception as e:
            logger.warning("%s failed (%s)", model_name, e)

    logger.warning("All Gemini image models failed, falling back to PIL")
    return False


def create_image(description: str, output_dir: str | None = None) -> str:
    """Generate an image: tries Gemini native generation first, falls back to PIL."""

    # Security: Sanitize output directory and filename
    if output_dir:
        output_path = Path(output_dir).resolve()
        # Verify output path is within allowed directory
        try:
            output_path.relative_to(ALLOWED_OUTPUT_DIR)
        except ValueError:
            raise ValueError(f"Output directory must be within {ALLOWED_OUTPUT_DIR}") from None  # noqa: B904
    else:
        output_path = ALLOWED_OUTPUT_DIR / "images"
        output_path.mkdir(parents=True, exist_ok=True)

    safe_title = sanitize_filename(description[:50])
    filepath = str(output_path / f"{safe_title}.png")

    # Try Gemini native image generation
    if _try_gemini_image_generation(description, filepath):
        return filepath

    # Fallback: PIL diagram with improved visuals
    import math

    from PIL import Image, ImageDraw, ImageFont

    # Extract a clean prompt for the diagram structure
    clean_desc = _extract_image_prompt(description)

    prompt = (
        f"Create a structured diagram for: '{clean_desc}'\n\n"
        f"Return JSON: {{"
        f'"title": "...", '
        f'"style": "org_chart"|"flowchart"|"architecture"|"simple", '
        f'"boxes": [{{"label": "...", "sublabel": "...", "x": 100, "y": 100, '
        f'"width": 180, "height": 70, "color": "#00D2FF"|"#FF6B6B"|"#4ECB71"|"#FFD93D"|"#A78BFA"}}, ...], '
        f'"connections": [{{"from": 0, "to": 1, "label": "...", "style": "solid"|"dashed"}}, ...]}}\n\n'
        f"Canvas: 1600x1000. Space boxes well. Use color to group related items. Only return JSON."
    )

    content = _ask_gemini(prompt, system=(
        "You generate professional diagram layouts as JSON. "
        "Use the full canvas. Assign distinct colors to different groups/departments. "
        "Add sublabels for roles/descriptions. If internal data is provided, use it accurately."))

    try:
        data = _parse_json(content)
    except json.JSONDecodeError:
        data = {
            "title": clean_desc[:60],
            "style": "simple",
            "boxes": [{"label": clean_desc[:40], "x": 600, "y": 400, "width": 400, "height": 200}],
            "connections": [],
        }

    width, height = 1600, 1000
    bg_color = (18, 18, 32)
    text_color = (255, 255, 255)
    muted_color = (160, 160, 200)

    img = Image.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(img)

    # Subtle grid background
    for gx in range(0, width, 40):
        draw.line([(gx, 0), (gx, height)], fill=(25, 25, 45), width=1)
    for gy in range(0, height, 40):
        draw.line([(0, gy), (width, gy)], fill=(25, 25, 45), width=1)

    # Load fonts
    try:
        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 36)
        box_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        sub_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 12)
        label_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 11)
    except OSError:
        title_font = box_font = sub_font = label_font = ImageFont.load_default()  # type: ignore[assignment]

    # Draw title with accent line
    title = data.get("title", "Diagram")
    title_bbox = draw.textbbox((0, 0), title, font=title_font)
    title_w = title_bbox[2] - title_bbox[0]
    draw.text(((width - title_w) / 2, 25), title, fill=(0, 210, 255), font=title_font)
    draw.line([(width/2 - title_w/2, 70), (width/2 + title_w/2, 70)], fill=(0, 210, 255), width=2)

    # Subtitle with timestamp
    sub_text = f"Generated by NEXUS  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    sub_bbox = draw.textbbox((0, 0), sub_text, font=sub_font)
    draw.text(((width - (sub_bbox[2] - sub_bbox[0])) / 2, 78), sub_text, fill=muted_color, font=sub_font)

    boxes = data.get("boxes", [])
    connections = data.get("connections", [])

    def hex_to_rgb(h):
        h = h.lstrip('#')
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def darken(rgb, factor=0.3):
        return tuple(max(0, int(c * factor)) for c in rgb)

    # Draw connections
    for conn in connections:
        fi, ti = conn.get("from", 0), conn.get("to", 0)
        if fi < len(boxes) and ti < len(boxes):
            fb, tb = boxes[fi], boxes[ti]
            fx = fb.get("x", 0) + fb.get("width", 100) / 2
            fy = fb.get("y", 0) + fb.get("height", 60) / 2
            tx = tb.get("x", 0) + tb.get("width", 100) / 2
            ty = tb.get("y", 0) + tb.get("height", 60) / 2

            line_color = (60, 60, 100)
            style = conn.get("style", "solid")

            if style == "dashed":
                # Dashed line
                length = math.sqrt((tx-fx)**2 + (ty-fy)**2)
                if length > 0:
                    dx, dy = (tx-fx)/length, (ty-fy)/length
                    dash_len, gap_len = 8, 6
                    pos = 0
                    while pos < length:
                        sx = fx + dx * pos
                        sy = fy + dy * pos
                        ex = fx + dx * min(pos + dash_len, length)
                        ey = fy + dy * min(pos + dash_len, length)
                        draw.line([(sx, sy), (ex, ey)], fill=line_color, width=2)
                        pos += dash_len + gap_len
            else:
                draw.line([(fx, fy), (tx, ty)], fill=line_color, width=2)

            # Arrowhead
            angle = math.atan2(ty - fy, tx - fx)
            sz = 10
            pts = [
                (tx, ty),
                (tx - sz * math.cos(angle - math.pi/6), ty - sz * math.sin(angle - math.pi/6)),
                (tx - sz * math.cos(angle + math.pi/6), ty - sz * math.sin(angle + math.pi/6)),
            ]
            draw.polygon(pts, fill=line_color)

            if conn.get("label"):
                mx, my = (fx + tx) / 2, (fy + ty) / 2
                lb = draw.textbbox((0, 0), conn["label"], font=label_font)
                lw = lb[2] - lb[0]
                # Background pill for label
                draw.rounded_rectangle(
                    [mx - lw/2 - 6, my - 10, mx + lw/2 + 6, my + 6],
                    radius=4, fill=(30, 30, 50))
                draw.text((mx - lw/2, my - 8), conn["label"], fill=muted_color, font=label_font)

    # Draw boxes with rounded corners and color accents
    for box in boxes:
        x, y = box.get("x", 0), box.get("y", 0)
        w, h = box.get("width", 180), box.get("height", 70)
        label = box.get("label", "")
        sublabel = box.get("sublabel", "")
        accent = hex_to_rgb(box.get("color", "#00D2FF"))
        bg = darken(accent, 0.15)

        # Rounded rectangle with accent top border
        draw.rounded_rectangle([x, y, x+w, y+h], radius=8, fill=bg, outline=(50, 50, 70), width=1)
        draw.rounded_rectangle([x, y, x+w, y+4], radius=2, fill=accent)

        # Label (centered)
        lb = draw.textbbox((0, 0), label, font=box_font)
        lw, lh = lb[2] - lb[0], lb[3] - lb[1]
        ty_offset = (h - lh) / 2 - (6 if sublabel else 0)
        draw.text((x + (w - lw) / 2, y + ty_offset), label, fill=text_color, font=box_font)

        # Sublabel
        if sublabel:
            sb = draw.textbbox((0, 0), sublabel, font=sub_font)
            sw = sb[2] - sb[0]
            draw.text((x + (w - sw) / 2, y + ty_offset + lh + 4), sublabel, fill=muted_color, font=sub_font)

    # Footer
    footer = "NEXUS Virtual Company"
    fb = draw.textbbox((0, 0), footer, font=sub_font)
    draw.text((width - (fb[2] - fb[0]) - 20, height - 25), footer, fill=(50, 50, 80), font=sub_font)

    img.save(filepath, quality=95)
    return filepath


# ============================================
# DISPATCHER
# ============================================

def detect_doc_request(message: str) -> dict | None:
    """Check if a message is asking for a document. Returns format info or None."""
    msg_lower = message.lower()

    format_map = {
        "docx": ["docx", "word doc", "word document", ".docx"],
        "pptx": ["pptx", "powerpoint", "slide deck", "slides", "presentation", ".pptx", "pitch deck"],
        "pdf": ["pdf", ".pdf"],
        "image": ["image", "diagram", "chart", "architecture diagram", "flowchart", "infographic"],
    }

    for fmt, keywords in format_map.items():
        for kw in keywords:
            if kw in msg_lower:
                return {"format": fmt}

    # Check for generic "document" or "report" — default to docx
    if any(w in msg_lower for w in ["document", "report", "memo", "letter", "write up", "writeup", "one-pager"]):
        return {"format": "docx"}

    return None


async def generate_document(message: str, doc_info: dict) -> dict:
    """Generate a document and return the filepath."""
    import asyncio

    fmt = doc_info["format"]

    # Gather internal NEXUS data relevant to the request
    internal_context = await _gather_internal_context(message)

    # Search the web if the request needs public info
    web_context = ""
    search_query = _needs_web_enrichment(message, bool(internal_context))
    if search_query:
        web_context = await _gather_web_context(search_query)

    # Build enriched request with all available context
    enriched_request = message
    context_parts = []
    if internal_context:
        context_parts.append(internal_context)
    if web_context:
        context_parts.append(web_context)
    if context_parts:
        enriched_request = (
            f"{message}\n\n"
            f"IMPORTANT: Use the following real data to generate this document. "
            f"Prefer internal data over web results when both are available. "
            f"Do NOT make up or hallucinate information — use exactly what is provided.\n\n"
            + "\n".join(context_parts)
        )

    # Extract a title from the message
    title = _ask_gemini(
        f"Extract a short document title (3-8 words) from this request: {message}\n\nReturn only the title, nothing else.",
        system="You extract concise titles. Return only the title text."
    ).strip().strip('"').strip("'")

    # Security: Use sanitized output directory
    output_dir = str(ALLOWED_OUTPUT_DIR)
    ALLOWED_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Run document creation in a thread pool (they're sync)
    loop = asyncio.get_event_loop()

    if fmt == "docx":
        filepath = await loop.run_in_executor(None, create_docx, title, enriched_request, output_dir)
    elif fmt == "pptx":
        filepath = await loop.run_in_executor(None, create_pptx, title, enriched_request, output_dir)
    elif fmt == "pdf":
        filepath = await loop.run_in_executor(None, create_pdf, title, enriched_request, output_dir)
    elif fmt == "image":
        filepath = await loop.run_in_executor(None, create_image, enriched_request, output_dir)
    else:
        return {"error": f"Unsupported format: {fmt}"}

    return {
        "filepath": filepath,
        "title": title,
        "format": fmt,
        "filename": os.path.basename(filepath),
    }
